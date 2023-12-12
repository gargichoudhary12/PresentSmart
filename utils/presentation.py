import io
import json
import os

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from urllib.parse import quote_plus
from dotenv import load_dotenv

dir_path = 'static/presentations'

load_dotenv()
API_KEY = "Nvj3QMFx5rlTMY7JIp6AnuKtFsLMxaOSXuU3zTvLRk2EWTe2cydXQVjg"


def parse_response(response):
    slides = response.split('\n\n')
    slides_content = []
    for slide in slides:
        lines = slide.split('\n')
        title_line = lines[0]
        if ': ' in title_line:
            title = title_line.split(': ', 1)[1]
        else:
            title = title_line
        content_lines = [line for line in lines[1:] if line != 'Content:']
        content = '\n'.join(content_lines)
        keyword_line = [line for line in lines if 'Keyword:' or 'Keywords:' in line][0]
        keyword = keyword_line.split(': ', 1)[1]
        slides_content.append({'title': title, 'content': content, 'keyword': keyword})
    return slides_content


def search_pexels_images(keyword):
    query = quote_plus(keyword.lower())
    print("Query:", query)
    PEXELS_API_URL = f'https://api.pexels.com/v1/search?query={query}&per_page=1'
    print("URL:", PEXELS_API_URL)
    headers = {
        'Authorization': API_KEY
    }
    response = requests.get(PEXELS_API_URL, headers=headers)
    print("Response Status Code:", response.status_code)
    print("Response Content:", response.text)
    data = json.loads(response.text)
    if 'photos' in data:
        if len(data['photos']) > 0:
            return data['photos'][0]['src']['medium']
    return None

def delete_first_two_slides(presentation):
    slide_ids = [1, 0]
    for slide_id in slide_ids:
        if slide_id < len(presentation.slides):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_id])


def create_ppt(slides_content, template_choice, presentation_title, presenter_name, insert_image):
    template_path = os.path.join(dir_path, f"{template_choice}.pptx")

    prs = Presentation(template_path)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = presentation_title


    subtitle = slide.placeholders[1]
    subtitle.text = f"Presented by {presenter_name}"

    if template_choice == 'dark':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Times New Roman'
                run.font.color.rgb = RGBColor(255, 165, 0)

    if template_choice == 'light':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Times New Roman'
                run.font.color.rgb = RGBColor(0, 0, 0)

    if template_choice == 'black&white':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.color.rgb = RGBColor(0, 0, 0)


    for slide_content in slides_content:
        slide = prs.slides.add_slide(content_slide_layout)

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = slide_content['title']
                if template_choice == 'dark':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 165, 0)
                if template_choice == 'light':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(0, 0, 0)

                if template_choice == 'black&white':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Arial'
                            run.font.color.rgb = RGBColor(255, 255, 255)

            elif placeholder.placeholder_format.type == 7:
                placeholder.text = slide_content['content']
                if template_choice == 'dark':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 255, 255)

                if template_choice == 'light':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(0, 0, 0)

                elif template_choice == 'black&white':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(0, 0, 0)

        if insert_image:
            image_url = search_pexels_images(slide_content['keyword'])
            print("Image URL:", image_url)  #
            if image_url is not None:

                image_data = requests.get(image_url).content

                image_stream = io.BytesIO(image_data)

                slide_width = Inches(25)
                slide_height = Inches(15)

                image_width = Inches(6)
                image_height = Inches(4)

                left = slide_width - image_width
                top = slide_height - image_height - Inches(4)

                slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)


    slide = prs.slides.add_slide(content_slide_layout)
    if template_choice == 'dark':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 165, 0)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 255, 255)

    if template_choice == 'light':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 20, 147)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    elif template_choice == 'black&white':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(255, 20, 147)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    else:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)
            elif placeholder.placeholder_format.type == 7:

                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    delete_first_two_slides(prs)


    prs.save(os.path.join('generated', 'generated_presentation.pptx'))